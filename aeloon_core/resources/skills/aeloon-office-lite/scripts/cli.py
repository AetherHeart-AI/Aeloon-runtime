#!/usr/bin/env python3
"""Fast, deterministic local I/O for simple PDF, DOCX, PPTX, and XLSX files."""

from __future__ import annotations

import argparse
import html
import importlib.metadata
import importlib.util
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import date, datetime
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

VALID_ACTIONS = ("preflight", "read", "write", "render", "validate")
SCHEMA = "aeloon-office-lite/v1"
TSINGHUA_INDEX = "https://pypi.tuna.tsinghua.edu.cn/simple"
SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm"}
WRITABLE_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx"}
LEGACY_SUFFIXES = {".doc", ".ppt", ".xls", ".wps", ".dps", ".et"}
DEPENDENCIES = {
    "pypdf": ("pypdf", "pypdf"),
    "pypdfium2": ("pypdfium2", "pypdfium2"),
    "python-docx": ("docx", "python-docx"),
    "python-pptx": ("pptx", "python-pptx"),
    "openpyxl": ("openpyxl", "openpyxl"),
    "reportlab": ("reportlab", "reportlab"),
}


class SkillError(ValueError):
    """A user-correctable Office Lite failure."""


class MissingDependencyError(SkillError):
    def __init__(self, packages: list[str]):
        unique = sorted(set(packages))
        self.packages = unique
        super().__init__(f"missing Python package(s): {', '.join(unique)}")


def install_hints(packages: list[str]) -> dict[str, str]:
    names = " ".join(sorted(set(packages)))
    return {
        "notice": "安装新包前先告知用户；中国网络环境默认使用清华 PyPI 镜像。",
        "pip": f"python -m pip install -i {TSINGHUA_INDEX} {names}",
        "uv": f"uv pip install --index-url {TSINGHUA_INDEX} {names}",
    }


def _package_state(module: str, distribution: str) -> dict[str, Any]:
    available = importlib.util.find_spec(module) is not None
    version = None
    if available:
        try:
            version = importlib.metadata.version(distribution)
        except importlib.metadata.PackageNotFoundError:
            pass
    return {"available": available, "version": version}


def preflight() -> dict[str, Any]:
    packages = {
        name: _package_state(module, distribution)
        for name, (module, distribution) in DEPENDENCIES.items()
    }
    missing = [name for name, state in packages.items() if not state["available"]]
    return {
        "schema": "aeloon-office-lite-preflight/v1",
        "python": sys.version.split()[0],
        "formats": sorted(SUPPORTED_SUFFIXES),
        "packages": packages,
        "libreoffice": _libreoffice_command(),
        "ready": not missing,
        "missing": missing,
        "install": install_hints(missing) if missing else None,
    }


def _require(module: str, distribution: str) -> None:
    if importlib.util.find_spec(module) is None:
        raise MissingDependencyError([distribution])


def _local_path(value: str, *, must_exist: bool = True) -> Path:
    parsed = urlparse(value)
    windows_drive = re.match(r"^[A-Za-z]:[\\/]", value)
    if (parsed.scheme and not windows_drive) or parsed.netloc:
        raise SkillError("只接受本地文件路径，拒绝 URI 或远程地址")
    path = Path(value).expanduser().resolve(strict=False)
    if must_exist and not path.is_file():
        raise SkillError(f"输入文件不存在: {path}")
    return path


def _input_path(value: str) -> Path:
    path = _local_path(value)
    suffix = path.suffix.lower()
    if suffix in LEGACY_SUFFIXES:
        raise SkillError(f"不支持旧格式 {suffix}；请先用 Office/WPS 另存为 DOCX、PPTX 或 XLSX")
    if suffix not in SUPPORTED_SUFFIXES:
        raise SkillError(f"不支持的文件类型 {suffix or '(无扩展名)'}")
    return path


def _output_path(value: str, *, overwrite: bool) -> Path:
    path = _local_path(value, must_exist=False)
    if path.suffix.lower() not in WRITABLE_SUFFIXES:
        raise SkillError("输出扩展名必须是 .pdf、.docx、.pptx 或 .xlsx")
    if path.exists() and not overwrite:
        raise SkillError(f"输出已存在；如需替换请传入 --overwrite: {path}")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _atomic_save(output: Path, suffix: str, save: Any) -> None:
    handle, temp_name = tempfile.mkstemp(
        prefix=f".{output.stem}-", suffix=suffix, dir=output.parent
    )
    os.close(handle)
    temporary = Path(temp_name)
    try:
        save(temporary)
        os.replace(temporary, output)
    finally:
        temporary.unlink(missing_ok=True)


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return str(value).replace("\r\n", "\n").replace("\r", "\n")


def _md_cell(value: Any) -> str:
    return _cell_text(value).replace("|", "\\|").replace("\n", "<br>")


def _markdown_table(rows: list[list[Any]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    lines = ["| " + " | ".join(_md_cell(cell) for cell in header) + " |"]
    lines.append("| " + " | ".join("---" for _ in range(width)) + " |")
    lines.extend("| " + " | ".join(_md_cell(cell) for cell in row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _read_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    _require("pypdf", "pypdf")
    from pypdf import PdfReader

    reader = PdfReader(path)
    chunks = [f"# {path.name}"]
    empty_pages: list[int] = []
    characters = 0
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        chunks.append(f"\n## 第 {index} 页\n")
        if text:
            chunks.append(text)
            characters += len(text)
        else:
            empty_pages.append(index)
            chunks.append("[本页没有可提取文本；请运行 render 后使用视觉能力读取页图]")
    warnings = []
    if empty_pages:
        warnings.append({"code": "pages_without_text", "pages": empty_pages})
    return "\n".join(chunks).rstrip() + "\n", {
        "pages": len(reader.pages),
        "characters": characters,
        "warnings": warnings,
    }


def _read_docx(path: Path) -> tuple[str, dict[str, Any]]:
    _require("docx", "python-docx")
    from docx import Document
    from docx.table import Table

    document = Document(path)
    chunks = [f"# {path.name}"]
    tables = 0
    for item in document.iter_inner_content():
        if isinstance(item, Table):
            rows = [[cell.text for cell in row.cells] for row in item.rows]
            table = _markdown_table(rows)
            if table:
                chunks.extend(["", table])
                tables += 1
            continue
        text = item.text.strip()
        if not text:
            continue
        style = (item.style.name if item.style is not None else "").lower()
        heading = re.match(r"heading\s*([1-6])", style)
        if heading:
            chunks.extend(["", f"{'#' * (int(heading.group(1)) + 1)} {text}"])
        elif "list" in style:
            chunks.append(f"- {text}")
        else:
            chunks.extend(["", text])
    text = "\n".join(chunks).rstrip() + "\n"
    return text, {"paragraphs": len(document.paragraphs), "tables": tables, "warnings": []}


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    return ""


def _read_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    _require("pptx", "python-pptx")
    from pptx import Presentation

    presentation = Presentation(path)
    chunks = [f"# {path.name}"]
    table_count = 0
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        chunks.extend(["", f"## 第 {index} 页" + (f"：{title}" if title else "")])
        for shape in slide.shapes:
            if shape is slide.shapes.title:
                continue
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                chunks.extend(["", _markdown_table(rows)])
                table_count += 1
                continue
            text = _shape_text(shape)
            if text:
                chunks.extend(["", text])
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, KeyError):
            notes = ""
        if notes:
            chunks.extend(["", "### 备注", notes])
    text = "\n".join(chunks).rstrip() + "\n"
    return text, {"slides": len(presentation.slides), "tables": table_count, "warnings": []}


def _read_xlsx(path: Path, *, max_rows: int, max_columns: int) -> tuple[str, dict[str, Any]]:
    _require("openpyxl", "openpyxl")
    from openpyxl import load_workbook

    keep_vba = path.suffix.lower() == ".xlsm"
    workbook = load_workbook(path, read_only=True, data_only=True, keep_vba=keep_vba)
    chunks = [f"# {path.name}"]
    warnings: list[dict[str, Any]] = []
    sheet_stats = []
    try:
        for sheet in workbook.worksheets:
            chunks.extend(["", f"## 工作表：{sheet.title}", ""])
            rows: list[list[Any]] = []
            for row_index, row in enumerate(
                sheet.iter_rows(max_row=max_rows, max_col=max_columns, values_only=True), start=1
            ):
                rows.append(list(row))
                if row_index >= max_rows:
                    break
            while rows and all(value is None for value in rows[-1]):
                rows.pop()
            if rows:
                chunks.append(_markdown_table(rows))
            else:
                chunks.append("[空工作表]")
            truncated = sheet.max_row > max_rows or sheet.max_column > max_columns
            if truncated:
                warnings.append(
                    {
                        "code": "sheet_truncated",
                        "sheet": sheet.title,
                        "source_rows": sheet.max_row,
                        "source_columns": sheet.max_column,
                    }
                )
            sheet_stats.append(
                {"name": sheet.title, "rows": sheet.max_row, "columns": sheet.max_column}
            )
    finally:
        workbook.close()
    return "\n".join(chunks).rstrip() + "\n", {"sheets": sheet_stats, "warnings": warnings}


def read_document(
    path: Path, *, max_rows: int = 200, max_columns: int = 50
) -> tuple[str, dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _read_xlsx(path, max_rows=max_rows, max_columns=max_columns)
    raise SkillError(f"不支持读取 {suffix}")


def _load_spec(path: Path) -> dict[str, Any]:
    try:
        spec = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise SkillError(f"无法读取 UTF-8 JSON spec: {exc}") from exc
    if not isinstance(spec, dict) or spec.get("schema") != SCHEMA:
        raise SkillError(f"spec.schema 必须是 {SCHEMA}")
    return spec


def _blocks(spec: dict[str, Any]) -> list[dict[str, Any]]:
    blocks = spec.get("blocks", [])
    if not isinstance(blocks, list) or any(not isinstance(block, dict) for block in blocks):
        raise SkillError("spec.blocks 必须是对象数组")
    return blocks


def _table_values(block: dict[str, Any]) -> list[list[Any]]:
    headers = block.get("headers", [])
    rows = block.get("rows", [])
    if not isinstance(headers, list) or not isinstance(rows, list):
        raise SkillError("table.headers 和 table.rows 必须是数组")
    if any(not isinstance(row, list) for row in rows):
        raise SkillError("table.rows 的每一行必须是数组")
    return ([headers] if headers else []) + rows


def _write_docx(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    _require("docx", "python-docx")
    from docx import Document
    from docx.shared import Inches, Pt

    document = Document()
    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(11)
    title = _cell_text(spec.get("title")).strip()
    if title:
        document.add_heading(title, level=0)
    for block in _blocks(spec):
        kind = block.get("type")
        if kind == "heading":
            level = max(1, min(6, int(block.get("level", 1))))
            document.add_heading(_cell_text(block.get("text")), level=level)
        elif kind == "paragraph":
            document.add_paragraph(_cell_text(block.get("text")))
        elif kind == "bullets":
            items = block.get("items", [])
            if not isinstance(items, list):
                raise SkillError("bullets.items 必须是数组")
            for item in items:
                document.add_paragraph(_cell_text(item), style="List Bullet")
        elif kind == "table":
            values = _table_values(block)
            if values:
                columns = max(len(row) for row in values)
                table = document.add_table(rows=len(values), cols=columns)
                table.style = "Table Grid"
                for row_index, row in enumerate(values):
                    for column_index, value in enumerate(row):
                        table.cell(row_index, column_index).text = _cell_text(value)
        elif kind == "pagebreak":
            document.add_page_break()
        elif kind == "image":
            image = _local_path(_cell_text(block.get("path")))
            width = float(block.get("width_inches", 5.0))
            document.add_picture(str(image), width=Inches(max(0.5, min(width, 7.0))))
        else:
            raise SkillError(f"不支持的 DOCX block.type: {kind}")
    core = document.core_properties
    core.title = title
    core.author = _cell_text(spec.get("author"))
    _atomic_save(output, ".docx", document.save)
    return {"blocks": len(_blocks(spec)), "title": title}


def _write_pptx(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    _require("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    title = _cell_text(spec.get("title")).strip()
    subtitle = _cell_text(spec.get("subtitle")).strip()
    if title:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
    slides = spec.get("slides", [])
    if not isinstance(slides, list) or any(not isinstance(slide, dict) for slide in slides):
        raise SkillError("spec.slides 必须是对象数组")
    for item in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12), Inches(0.65))
        title_frame = title_box.text_frame
        title_frame.text = _cell_text(item.get("title"))
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        if "table" in item:
            table_spec = item["table"]
            if not isinstance(table_spec, dict):
                raise SkillError("slide.table 必须是对象")
            values = _table_values(table_spec)
            if values:
                columns = max(len(row) for row in values)
                shape = slide.shapes.add_table(
                    len(values), columns, Inches(0.75), Inches(1.35), Inches(11.8), Inches(5.3)
                )
                table = shape.table
                for row_index, row in enumerate(values):
                    for column_index, value in enumerate(row):
                        cell = table.cell(row_index, column_index)
                        cell.text = _cell_text(value)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(15)
                            if row_index == 0:
                                paragraph.font.bold = True
        elif "image" in item:
            image = _local_path(_cell_text(item.get("image")))
            slide.shapes.add_picture(
                str(image), Inches(1.2), Inches(1.35), width=Inches(10.9), height=Inches(5.6)
            )
        else:
            body = slide.shapes.add_textbox(Inches(0.95), Inches(1.35), Inches(11.3), Inches(5.3))
            frame = body.text_frame
            frame.word_wrap = True
            bullets = item.get("bullets")
            if bullets is not None:
                if not isinstance(bullets, list):
                    raise SkillError("slide.bullets 必须是数组")
                values = bullets or [""]
                for index, value in enumerate(values):
                    paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    paragraph.text = _cell_text(value)
                    paragraph.level = 0
                    paragraph.font.size = Pt(24)
                    paragraph.space_after = Pt(14)
            else:
                frame.text = _cell_text(item.get("text"))
                frame.paragraphs[0].font.size = Pt(24)
                frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        notes = _cell_text(item.get("notes")).strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    _atomic_save(output, ".pptx", presentation.save)
    return {"slides": len(presentation.slides), "content_slides": len(slides), "title": title}


def _safe_sheet_name(value: Any, existing: set[str]) -> str:
    name = re.sub(r"[\\/*?:\[\]]", "_", _cell_text(value).strip())[:31] or "Sheet"
    base = name
    index = 2
    while name in existing:
        suffix = f"-{index}"
        name = base[: 31 - len(suffix)] + suffix
        index += 1
    return name


def _write_xlsx(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    _require("openpyxl", "openpyxl")
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    workbook.remove(workbook.active)
    sheets = spec.get("sheets", [])
    if (
        not isinstance(sheets, list)
        or not sheets
        or any(not isinstance(sheet, dict) for sheet in sheets)
    ):
        raise SkillError("spec.sheets 必须是非空对象数组")
    existing: set[str] = set()
    row_count = 0
    for sheet_spec in sheets:
        name = _safe_sheet_name(sheet_spec.get("name"), existing)
        existing.add(name)
        sheet = workbook.create_sheet(name)
        rows = sheet_spec.get("rows", [])
        if not isinstance(rows, list) or any(not isinstance(row, list) for row in rows):
            raise SkillError("sheet.rows 必须是二维数组")
        for row in rows:
            sheet.append(row)
        row_count += len(rows)
        if rows and sheet_spec.get("header", True):
            fill = PatternFill("solid", fgColor="D9EAF7")
            for cell in sheet[1]:
                cell.font = Font(bold=True)
                cell.fill = fill
        freeze = sheet_spec.get("freeze")
        if freeze:
            sheet.freeze_panes = _cell_text(freeze)
        if rows and sheet_spec.get("auto_filter", False):
            sheet.auto_filter.ref = sheet.dimensions
        width = min(max((len(row) for row in rows), default=0), 100)
        for column_index in range(1, width + 1):
            values = [_cell_text(row[column_index - 1]) for row in rows if len(row) >= column_index]
            sheet.column_dimensions[get_column_letter(column_index)].width = min(
                max((len(value) for value in values), default=8) + 2, 40
            )
    _atomic_save(output, output.suffix.lower(), workbook.save)
    return {"sheets": len(sheets), "rows": row_count}


def _pdf_styles():
    _require("reportlab", "reportlab")
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    try:
        pdfmetrics.getFont("STSong-Light")
    except KeyError:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    styles = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "LiteTitle",
            parent=styles["Title"],
            fontName="STSong-Light",
            fontSize=22,
            leading=30,
            spaceAfter=18,
        ),
        "heading": ParagraphStyle(
            "LiteHeading",
            parent=styles["Heading1"],
            fontName="STSong-Light",
            fontSize=16,
            leading=22,
            spaceBefore=12,
            spaceAfter=8,
        ),
        "body": ParagraphStyle(
            "LiteBody",
            parent=styles["BodyText"],
            fontName="STSong-Light",
            fontSize=11,
            leading=17,
            spaceAfter=8,
        ),
    }


def _pdf_text(value: Any) -> str:
    return html.escape(_cell_text(value)).replace("\n", "<br/>")


def _write_pdf(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    _require("reportlab", "reportlab")
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    styles = _pdf_styles()
    story = []
    title = _cell_text(spec.get("title")).strip()
    if title:
        story.append(Paragraph(_pdf_text(title), styles["title"]))
    for block in _blocks(spec):
        kind = block.get("type")
        if kind == "heading":
            story.append(Paragraph(_pdf_text(block.get("text")), styles["heading"]))
        elif kind == "paragraph":
            story.append(Paragraph(_pdf_text(block.get("text")), styles["body"]))
        elif kind == "bullets":
            items = block.get("items", [])
            if not isinstance(items, list):
                raise SkillError("bullets.items 必须是数组")
            for item in items:
                story.append(Paragraph("• " + _pdf_text(item), styles["body"]))
        elif kind == "table":
            values = _table_values(block)
            if values:
                data = [
                    [Paragraph(_pdf_text(value), styles["body"]) for value in row] for row in values
                ]
                table = Table(data, repeatRows=1, hAlign="LEFT")
                commands = [
                    ("FONTNAME", (0, 0), (-1, -1), "STSong-Light"),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ]
                table.setStyle(TableStyle(commands))
                story.extend([table, Spacer(1, 4 * mm)])
        elif kind == "pagebreak":
            story.append(PageBreak())
        else:
            raise SkillError(f"不支持的 PDF block.type: {kind}")

    def save(path: Path) -> None:
        document = SimpleDocTemplate(
            str(path),
            pagesize=A4,
            rightMargin=18 * mm,
            leftMargin=18 * mm,
            topMargin=16 * mm,
            bottomMargin=16 * mm,
            title=title,
            author=_cell_text(spec.get("author")),
        )
        document.build(story)

    _atomic_save(output, ".pdf", save)
    return {"blocks": len(_blocks(spec)), "title": title}


def write_document(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    suffix = output.suffix.lower()
    if suffix == ".pdf":
        return _write_pdf(spec, output)
    if suffix == ".docx":
        return _write_docx(spec, output)
    if suffix == ".pptx":
        return _write_pptx(spec, output)
    if suffix == ".xlsx":
        return _write_xlsx(spec, output)
    raise SkillError(f"不支持写入 {suffix}")


def validate_document(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    details: dict[str, Any]
    warnings: list[dict[str, Any]] = []
    if suffix == ".pdf":
        _require("pypdf", "pypdf")
        from pypdf import PdfReader

        reader = PdfReader(path)
        details = {"pages": len(reader.pages), "encrypted": reader.is_encrypted}
        if not reader.pages:
            warnings.append({"code": "empty_document"})
    elif suffix == ".docx":
        _require("docx", "python-docx")
        from docx import Document

        document = Document(path)
        details = {"paragraphs": len(document.paragraphs), "tables": len(document.tables)}
    elif suffix == ".pptx":
        _require("pptx", "python-pptx")
        from pptx import Presentation

        presentation = Presentation(path)
        details = {"slides": len(presentation.slides)}
        if not presentation.slides:
            warnings.append({"code": "empty_document"})
    elif suffix in {".xlsx", ".xlsm"}:
        _require("openpyxl", "openpyxl")
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=False, keep_vba=suffix == ".xlsm")
        try:
            details = {"sheets": len(workbook.sheetnames), "sheet_names": workbook.sheetnames}
        finally:
            workbook.close()
    else:
        raise SkillError(f"不支持验证 {suffix}")
    return {
        "schema": "aeloon-office-lite-validation/v1",
        "file": str(path),
        "format": suffix.lstrip("."),
        "bytes": path.stat().st_size,
        "valid": True,
        "details": details,
        "warnings": warnings,
    }


def _libreoffice_command() -> str | None:
    discovered = shutil.which("libreoffice") or shutil.which("soffice")
    if discovered:
        return discovered
    mac = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    return str(mac) if mac.is_file() else None


def _render_pdf(source: Path, output_dir: Path, *, dpi: int, overwrite: bool) -> list[Path]:
    _require("pypdfium2", "pypdfium2")
    import pypdfium2

    output_dir.mkdir(parents=True, exist_ok=True)
    document = pypdfium2.PdfDocument(source)
    outputs: list[Path] = []
    try:
        digits = max(3, len(str(len(document))))
        for index in range(len(document)):
            output = output_dir / f"page-{index + 1:0{digits}d}.png"
            if output.exists() and not overwrite:
                raise SkillError(f"渲染文件已存在；如需替换请传入 --overwrite: {output}")
            page = document[index]
            try:
                bitmap = page.render(scale=dpi / 72)
                image = bitmap.to_pil()
                image.save(output, format="PNG", optimize=False)
                outputs.append(output)
            finally:
                page.close()
    finally:
        document.close()
    return outputs


def render_document(source: Path, output_dir: Path, *, dpi: int, overwrite: bool) -> list[Path]:
    if source.suffix.lower() == ".pdf":
        return _render_pdf(source, output_dir, dpi=dpi, overwrite=overwrite)
    office = _libreoffice_command()
    if office is None:
        raise SkillError(
            "渲染 DOCX/PPTX/XLSX 需要 LibreOffice；macOS: brew install --cask libreoffice；"
            "Ubuntu: sudo apt-get install libreoffice"
        )
    with tempfile.TemporaryDirectory(prefix="aeloon-office-lite-") as temporary:
        command = [office, "--headless", "--convert-to", "pdf", "--outdir", temporary, str(source)]
        result = subprocess.run(command, capture_output=True, text=True, timeout=120, check=False)
        pdf = Path(temporary) / f"{source.stem}.pdf"
        if result.returncode != 0 or not pdf.is_file():
            message = (result.stderr or result.stdout).strip()
            raise SkillError(f"LibreOffice 转换失败: {message or '未生成 PDF'}")
        return _render_pdf(pdf, output_dir, dpi=dpi, overwrite=overwrite)


def _json_dump(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, indent=2) + "\n"


def _build_parser(action: str) -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog=f"aeloon-office-lite {action}")
    if action == "preflight":
        return parser
    if action == "read":
        parser.add_argument("input")
        parser.add_argument("--output")
        parser.add_argument("--metadata")
        parser.add_argument("--max-rows", type=int, default=200)
        parser.add_argument("--max-columns", type=int, default=50)
        return parser
    if action == "write":
        parser.add_argument("spec")
        parser.add_argument("output")
        parser.add_argument("--overwrite", action="store_true")
        return parser
    if action == "validate":
        parser.add_argument("input")
        parser.add_argument("--output")
        return parser
    if action == "render":
        parser.add_argument("input")
        parser.add_argument("--output-dir", required=True)
        parser.add_argument("--dpi", type=int, default=120)
        parser.add_argument("--overwrite", action="store_true")
        return parser
    raise SkillError(f"未知 action: {action}")


def main(argv: list[str] | None = None) -> int:
    arguments = list(sys.argv[1:] if argv is None else argv)
    if not arguments or arguments[0] not in VALID_ACTIONS:
        print(
            _json_dump({"error": "unknown_action", "valid_actions": list(VALID_ACTIONS)}),
            file=sys.stderr,
            end="",
        )
        return 2
    action = arguments.pop(0)
    try:
        args = _build_parser(action).parse_args(arguments)
        if action == "preflight":
            print(_json_dump(preflight()), end="")
        elif action == "read":
            source = _input_path(args.input)
            if args.max_rows < 1 or args.max_columns < 1:
                raise SkillError("--max-rows 和 --max-columns 必须大于 0")
            markdown, metadata = read_document(
                source, max_rows=args.max_rows, max_columns=args.max_columns
            )
            if args.output:
                output = _local_path(args.output, must_exist=False)
                if output.exists():
                    raise SkillError(f"输出已存在: {output}")
                output.parent.mkdir(parents=True, exist_ok=True)
                output.write_text(markdown, encoding="utf-8")
            else:
                print(markdown, end="")
            if args.metadata:
                meta_output = _local_path(args.metadata, must_exist=False)
                if meta_output.exists():
                    raise SkillError(f"metadata 输出已存在: {meta_output}")
                meta_output.parent.mkdir(parents=True, exist_ok=True)
                meta_output.write_text(_json_dump(metadata), encoding="utf-8")
        elif action == "write":
            spec_path = _local_path(args.spec)
            output = _output_path(args.output, overwrite=args.overwrite)
            result = write_document(_load_spec(spec_path), output)
            print(
                _json_dump(
                    {"schema": "aeloon-office-lite-write/v1", "output": str(output), **result}
                ),
                end="",
            )
        elif action == "validate":
            result = validate_document(_input_path(args.input))
            payload = _json_dump(result)
            if args.output:
                output = _local_path(args.output, must_exist=False)
                if output.exists():
                    raise SkillError(f"输出已存在: {output}")
                output.parent.mkdir(parents=True, exist_ok=True)
                output.write_text(payload, encoding="utf-8")
            else:
                print(payload, end="")
        elif action == "render":
            if not 36 <= args.dpi <= 300:
                raise SkillError("--dpi 必须在 36 到 300 之间")
            source = _input_path(args.input)
            output_dir = _local_path(args.output_dir, must_exist=False)
            outputs = render_document(source, output_dir, dpi=args.dpi, overwrite=args.overwrite)
            print(
                _json_dump(
                    {
                        "schema": "aeloon-office-lite-render/v1",
                        "pages": len(outputs),
                        "files": [str(path) for path in outputs],
                    }
                ),
                end="",
            )
        return 0
    except MissingDependencyError as exc:
        print(
            _json_dump(
                {
                    "error": "missing_dependency",
                    "message": str(exc),
                    "install": install_hints(exc.packages),
                }
            ),
            file=sys.stderr,
            end="",
        )
        return 3
    except (SkillError, OSError, ValueError) as exc:
        print(
            _json_dump({"error": "office_lite_error", "message": str(exc)}), file=sys.stderr, end=""
        )
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
