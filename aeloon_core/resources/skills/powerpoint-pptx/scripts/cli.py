#!/usr/bin/env python3
"""Pure-Python PowerPoint build, template editing, rendering, and QA CLI."""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import urlparse
from zipfile import BadZipFile, ZipFile

from markdown_it import MarkdownIt
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

EMU_PER_INCH = 914400
VALID_ACTIONS = ("build", "inspect-template", "apply-template", "validate", "render")
PLACEHOLDER_RE = re.compile(
    r"(?:\b(?:lorem|ipsum|placeholder|tbd|xxxx+)\b|待(?:填写|补充)|示例(?:文字|标题))",
    re.IGNORECASE,
)
ADDRESS_RE = re.compile(
    r"^slide/(?P<slide>\d+)/shape/(?P<shape>\d+)"
    r"(?:/paragraph/(?P<paragraph>\d+)/run/(?P<run>\d+))?$"
)
CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}
THEME = {
    "background": "F8FAFC",
    "surface": "FFFFFF",
    "ink": "172033",
    "muted": "526078",
    "accent": "2878D0",
    "accent_soft": "E7F1FC",
    "border": "DCE3EC",
}


class SkillError(RuntimeError):
    """A user-correctable CLI error."""


@dataclass
class SlideContent:
    title: str = ""
    paragraphs: list[str] = field(default_factory=list)
    bullets: list[dict[str, Any]] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    images: list[dict[str, str]] = field(default_factory=list)
    charts: list[dict[str, Any]] = field(default_factory=list)
    notes: str = ""


def _rgb(value: str) -> RGBColor:
    value = value.removeprefix("#")
    if not re.fullmatch(r"[0-9a-fA-F]{6}", value):
        raise SkillError(f"invalid six-digit color: {value!r}")
    return RGBColor.from_string(value.upper())


def _local_file(value: str | Path, *, suffixes: set[str] | None = None) -> Path:
    raw = str(value)
    parsed = urlparse(raw)
    if parsed.scheme or parsed.netloc:
        raise SkillError("only local files are supported; URI input was rejected")
    path = Path(raw).expanduser().resolve()
    if not path.is_file():
        raise SkillError(f"input file does not exist: {path}")
    if suffixes and path.suffix.lower() not in suffixes:
        allowed = ", ".join(sorted(suffixes))
        raise SkillError(f"unsupported input extension {path.suffix!r}; expected {allowed}")
    return path


def _output_file(
    value: str | Path,
    *,
    source: Path | None = None,
    overwrite: bool = False,
    suffix: str | None = None,
) -> Path:
    path = Path(value).expanduser().resolve()
    if source is not None and path == source.resolve():
        raise SkillError("input and output paths must differ; the input is never overwritten")
    if suffix and path.suffix.lower() != suffix:
        raise SkillError(f"output must use the {suffix} extension")
    if path.exists() and not overwrite:
        raise SkillError(f"output already exists: {path}; pass --overwrite to replace it")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _write_json(path: Path | None, payload: dict[str, Any], *, overwrite: bool) -> None:
    if path is None:
        return
    target = _output_file(path, overwrite=overwrite, suffix=".json")
    target.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _save_presentation(prs: Presentation, output: Path) -> None:
    with tempfile.NamedTemporaryFile(
        prefix=f".{output.stem}-", suffix=".pptx", dir=output.parent, delete=False
    ) as handle:
        temporary = Path(handle.name)
    try:
        prs.save(temporary)
        os.replace(temporary, output)
    finally:
        temporary.unlink(missing_ok=True)


def _inline_content(token: Any) -> tuple[str, list[dict[str, str]]]:
    text: list[str] = []
    images: list[dict[str, str]] = []
    for child in token.children or []:
        if child.type in {"text", "code_inline", "html_inline"}:
            text.append(child.content)
        elif child.type in {"softbreak", "hardbreak"}:
            text.append("\n")
        elif child.type == "image":
            images.append({"src": child.attrGet("src") or "", "alt": child.content or ""})
        elif child.type == "link_open":
            continue
    return "".join(text).strip(), images


def _validate_chart(data: Any) -> dict[str, Any]:
    if not isinstance(data, dict):
        raise SkillError("chart block must contain a JSON object")
    chart_type = str(data.get("type", "column")).lower()
    if chart_type not in CHART_TYPES:
        raise SkillError(f"unsupported chart type {chart_type!r}")
    categories = data.get("categories")
    series = data.get("series")
    if not isinstance(categories, list) or not categories:
        raise SkillError("chart categories must be a non-empty array")
    if not isinstance(series, list) or not series:
        raise SkillError("chart series must be a non-empty array")
    normalized_series = []
    for item in series:
        if not isinstance(item, dict) or not str(item.get("name", "")).strip():
            raise SkillError("each chart series requires a non-empty name")
        values = item.get("values")
        if not isinstance(values, list) or len(values) != len(categories):
            raise SkillError("each chart series must have one value per category")
        if any(not isinstance(value, (int, float)) for value in values):
            raise SkillError("chart values must be numeric")
        normalized_series.append({"name": str(item["name"]), "values": values})
    return {
        "type": chart_type,
        "categories": [str(category) for category in categories],
        "series": normalized_series,
    }


def _parse_slide(markdown: str) -> SlideContent:
    parser = MarkdownIt("commonmark").enable("table")
    tokens = parser.parse(markdown)
    slide = SlideContent()
    list_stack: list[bool] = []
    in_list_item = False
    table_rows: list[list[str]] | None = None
    table_row: list[str] | None = None
    in_table_cell = False

    for index, token in enumerate(tokens):
        if token.type == "heading_open":
            following = tokens[index + 1] if index + 1 < len(tokens) else None
            if following is not None and following.type == "inline":
                text, images = _inline_content(following)
                if not slide.title:
                    slide.title = text
                elif text:
                    slide.paragraphs.append(text)
                slide.images.extend(images)
        elif token.type in {"bullet_list_open", "ordered_list_open"}:
            list_stack.append(token.type == "ordered_list_open")
        elif token.type in {"bullet_list_close", "ordered_list_close"}:
            if list_stack:
                list_stack.pop()
        elif token.type == "list_item_open":
            in_list_item = True
        elif token.type == "list_item_close":
            in_list_item = False
        elif token.type == "table_open":
            table_rows = []
        elif token.type == "table_close":
            if table_rows:
                slide.tables.append(table_rows)
            table_rows = None
        elif token.type == "tr_open":
            table_row = []
        elif token.type == "tr_close":
            if table_rows is not None and table_row is not None:
                table_rows.append(table_row)
            table_row = None
        elif token.type in {"th_open", "td_open"}:
            in_table_cell = True
        elif token.type in {"th_close", "td_close"}:
            in_table_cell = False
        elif token.type == "inline":
            previous = tokens[index - 1].type if index else ""
            if previous == "heading_open":
                continue
            text, images = _inline_content(token)
            slide.images.extend(images)
            if in_table_cell and table_row is not None:
                table_row.append(text)
            elif in_list_item and text:
                slide.bullets.append(
                    {
                        "text": text,
                        "level": max(0, len(list_stack) - 1),
                        "ordered": bool(list_stack and list_stack[-1]),
                    }
                )
            elif text:
                slide.paragraphs.append(text)
        elif token.type == "fence":
            info = token.info.strip().lower()
            if info == "notes":
                slide.notes = token.content.strip()
            elif info == "chart":
                try:
                    chart_data = json.loads(token.content)
                except json.JSONDecodeError as exc:
                    raise SkillError(f"invalid chart JSON: {exc}") from exc
                slide.charts.append(_validate_chart(chart_data))
    if not slide.title.strip():
        raise SkillError("every slide requires a heading to use as its title")
    return slide


def parse_markdown(source: Path) -> list[SlideContent]:
    markdown = source.read_text(encoding="utf-8")
    chunks = re.split(r"(?m)^\s*---\s*$", markdown)
    slides = [_parse_slide(chunk) for chunk in chunks if chunk.strip()]
    if not slides:
        raise SkillError("Markdown input did not contain any slides")
    for slide in slides:
        for image in slide.images:
            image_path = Path(image["src"]).expanduser()
            if not image_path.is_absolute():
                image_path = source.parent / image_path
            image["src"] = str(image_path.resolve())
            if not image_path.is_file():
                raise SkillError(f"image does not exist: {image_path}")
    return slides


def _set_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_text_box(
    slide: Any,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str,
    size: float,
    color: str,
    bold: bool = False,
    name: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _add_footer(slide: Any, page: int, footer: str, *, font: str) -> None:
    _add_text_box(
        slide,
        footer,
        x=0.72,
        y=7.08,
        w=10.8,
        h=0.18,
        font=font,
        size=9,
        color=THEME["muted"],
        name="AeloonFooter",
    )
    _add_text_box(
        slide,
        f"{page:02d}",
        x=12.0,
        y=7.02,
        w=0.6,
        h=0.24,
        font=font,
        size=10,
        color=THEME["accent"],
        bold=True,
        name="AeloonPageNumber",
        align=PP_ALIGN.RIGHT,
    )


def _add_cover(slide: Any, content: SlideContent, *, font: str, accent: str) -> None:
    _set_background(slide, THEME["background"])
    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.92), Inches(0), Inches(3.41), Inches(7.5)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = _rgb(accent)
    block.line.fill.background()
    marker = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.28), Inches(0.12), Inches(0.56)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = _rgb(accent)
    marker.line.fill.background()
    _add_text_box(
        slide,
        content.title,
        x=1.05,
        y=1.2,
        w=8.1,
        h=2.35,
        font=font,
        size=34,
        color=THEME["ink"],
        bold=True,
        name="AeloonTitle",
        valign=MSO_ANCHOR.MIDDLE,
    )
    subtitle = content.paragraphs[0] if content.paragraphs else ""
    if subtitle:
        _add_text_box(
            slide,
            subtitle,
            x=1.05,
            y=3.75,
            w=7.7,
            h=1.2,
            font=font,
            size=18,
            color=THEME["muted"],
            name="AeloonSubtitle",
        )
    brand = _add_text_box(
        slide,
        "AELOON",
        x=10.46,
        y=6.65,
        w=2.2,
        h=0.35,
        font=font,
        size=11,
        color=THEME["surface"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    brand.fill.solid()
    brand.fill.fore_color.rgb = _rgb(accent)
    brand.line.fill.background()


def _set_bullet(paragraph: Any, *, ordered: bool, start: int = 1) -> None:
    properties = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        node = properties.find(tag, properties.nsmap)
        if node is not None:
            properties.remove(node)
    if ordered:
        bullet = OxmlElement("a:buAutoNum")
        bullet.set("type", "arabicPeriod")
        bullet.set("startAt", str(start))
    else:
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
    properties.insert(0, bullet)


def _add_body_text(
    slide: Any,
    content: SlideContent,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str,
) -> Any | None:
    entries: list[tuple[str, dict[str, Any] | None]] = [
        (paragraph, None) for paragraph in content.paragraphs
    ]
    entries.extend((str(item["text"]), item) for item in content.bullets)
    if not entries:
        return None
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = "AeloonBody"
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    for index, (text, bullet) in enumerate(entries):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.1
        if bullet:
            paragraph.level = min(int(bullet["level"]), 4)
            _set_bullet(paragraph, ordered=bool(bullet["ordered"]), start=index + 1)
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(17 if bullet else 16)
        run.font.color.rgb = _rgb(THEME["ink"])
    return shape


def _add_image(slide: Any, data: dict[str, str], *, x: float, y: float, w: float, h: float) -> Any:
    source = Path(data["src"])
    with Image.open(source) as image:
        pixel_width, pixel_height = image.size
    ratio = pixel_width / max(pixel_height, 1)
    box_ratio = w / h
    if ratio > box_ratio:
        image_w = w
        image_h = w / ratio
    else:
        image_h = h
        image_w = h * ratio
    shape = slide.shapes.add_picture(
        str(source),
        Inches(x + (w - image_w) / 2),
        Inches(y + (h - image_h) / 2),
        Inches(image_w),
        Inches(image_h),
    )
    shape.name = "AeloonImage"
    return shape


def _add_table(
    slide: Any,
    rows: list[list[str]],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str,
) -> Any:
    column_count = max((len(row) for row in rows), default=0)
    if not rows or not column_count:
        raise SkillError("table must contain at least one cell")
    shape = slide.shapes.add_table(
        len(rows), column_count, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = "AeloonTable"
    table = shape.table
    for row_index, row in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                THEME["accent_soft"] if row_index == 0 else THEME["surface"]
            )
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font
                    run.font.size = Pt(13)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = _rgb(THEME["ink"])
    return shape


def _chart_data(data: dict[str, Any]) -> ChartData:
    chart_data = ChartData()
    chart_data.categories = data["categories"]
    for series in data["series"]:
        chart_data.add_series(series["name"], series["values"])
    return chart_data


def _add_chart(
    slide: Any,
    data: dict[str, Any],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    font: str,
) -> Any:
    shape = slide.shapes.add_chart(
        CHART_TYPES[data["type"]], Inches(x), Inches(y), Inches(w), Inches(h), _chart_data(data)
    )
    shape.name = "AeloonChart"
    chart = shape.chart
    chart.has_legend = len(data["series"]) > 1 or data["type"] in {"pie", "doughnut"}
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = font
        chart.legend.font.size = Pt(11)
    chart.has_title = False
    return shape


def _add_content_slide(
    slide: Any,
    content: SlideContent,
    *,
    page: int,
    font: str,
    footer: str,
    accent: str,
) -> None:
    _set_background(slide, THEME["background"])
    accent_block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(0.64),
        Inches(0.22),
        Inches(0.68),
    )
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = _rgb(accent)
    accent_block.line.fill.background()
    _add_text_box(
        slide,
        content.title,
        x=1.12,
        y=0.58,
        w=11.35,
        h=0.82,
        font=font,
        size=27,
        color=THEME["ink"],
        bold=True,
        name="AeloonTitle",
        valign=MSO_ANCHOR.MIDDLE,
    )
    content_y = 1.62
    content_h = 5.13
    has_visual = bool(content.images or content.charts or content.tables)
    if content.charts:
        _add_body_text(slide, content, x=0.8, y=content_y, w=3.6, h=content_h, font=font)
        _add_chart(
            slide,
            content.charts[0],
            x=4.65,
            y=content_y,
            w=7.85,
            h=content_h,
            font=font,
        )
    elif content.tables:
        body_height = 1.1 if (content.paragraphs or content.bullets) else 0
        if body_height:
            _add_body_text(slide, content, x=0.8, y=content_y, w=11.7, h=body_height, font=font)
        table_y = content_y + body_height + (0.18 if body_height else 0)
        _add_table(
            slide,
            content.tables[0],
            x=0.8,
            y=table_y,
            w=11.7,
            h=content_h - body_height,
            font=font,
        )
    elif content.images:
        _add_body_text(slide, content, x=0.8, y=content_y, w=5.45, h=content_h, font=font)
        _add_image(slide, content.images[0], x=6.58, y=content_y, w=5.92, h=content_h)
    else:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.8),
            Inches(content_y),
            Inches(11.7),
            Inches(content_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(THEME["surface"])
        card.line.color.rgb = _rgb(THEME["border"])
        _add_body_text(
            slide, content, x=1.12, y=content_y + 0.35, w=10.95, h=content_h - 0.7, font=font
        )
    if has_visual and len(content.images) + len(content.charts) + len(content.tables) > 1:
        _add_text_box(
            slide,
            "Additional visual blocks were omitted; split dense content into one idea per slide.",
            x=0.8,
            y=6.79,
            w=10.4,
            h=0.18,
            font=font,
            size=8,
            color=THEME["muted"],
            name="AeloonBuildWarning",
        )
    _add_footer(slide, page, footer, font=font)


def _set_notes(slide: Any, notes: str) -> None:
    if not notes:
        return
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes


def build_deck(
    source: Path,
    output: Path,
    *,
    title: str | None,
    author: str | None,
    accent: str,
    font: str,
    footer: str,
) -> dict[str, Any]:
    slides = parse_markdown(source)
    accent = accent.removeprefix("#").upper()
    _rgb(accent)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title or slides[0].title
    if author:
        prs.core_properties.author = author
    blank = prs.slide_layouts[6]
    for index, content in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        if index == 0:
            _add_cover(slide, content, font=font, accent=accent)
        else:
            _add_content_slide(
                slide,
                content,
                page=index + 1,
                font=font,
                footer=footer,
                accent=accent,
            )
        _set_notes(slide, content.notes)
    _save_presentation(prs, output)
    return {
        "schema": "ppt-build-result/v1",
        "source": str(source),
        "output": str(output),
        "slides": len(slides),
        "theme": "aeloon-minimal-16x9/v1",
    }


def _font_size(run: Any) -> float | None:
    return round(run.font.size.pt, 2) if run.font.size is not None else None


def _shape_text_detail(shape: Any, address: str) -> list[dict[str, Any]]:
    if not shape.has_text_frame:
        return []
    paragraphs = []
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
        paragraph_address = f"{address}/paragraph/{paragraph_index}"
        runs = []
        for run_index, run in enumerate(paragraph.runs, start=1):
            runs.append(
                {
                    "address": f"{paragraph_address}/run/{run_index}",
                    "text": run.text,
                    "font": {
                        "name": run.font.name,
                        "size_pt": _font_size(run),
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                    },
                }
            )
        paragraphs.append(
            {
                "address": paragraph_address,
                "level": paragraph.level,
                "text": paragraph.text,
                "runs": runs,
            }
        )
    return paragraphs


def _chart_detail(shape: Any) -> dict[str, Any] | None:
    if not shape.has_chart:
        return None
    chart = shape.chart
    result: dict[str, Any] = {"series": []}
    try:
        result["categories"] = [str(category.label) for category in chart.plots[0].categories]
    except (AttributeError, IndexError, TypeError):
        result["categories"] = []
    for series in chart.series:
        values = []
        try:
            values = list(series.values)
        except (AttributeError, TypeError):
            pass
        result["series"].append({"name": series.name, "values": values})
    return result


def inspect_template(source: Path) -> dict[str, Any]:
    prs = Presentation(source)
    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_address = f"slide/{slide_index}"
        shapes = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            address = f"{slide_address}/shape/{shape_index}"
            width_inches = shape.width / EMU_PER_INCH
            height_inches = shape.height / EMU_PER_INCH
            shapes.append(
                {
                    "address": address,
                    "name": shape.name,
                    "shape_type": int(shape.shape_type),
                    "placeholder_type": (
                        str(shape.placeholder_format.type) if shape.is_placeholder else None
                    ),
                    "position": {
                        "x_inches": round(shape.left / EMU_PER_INCH, 3),
                        "y_inches": round(shape.top / EMU_PER_INCH, 3),
                        "width_inches": round(width_inches, 3),
                        "height_inches": round(height_inches, 3),
                    },
                    "text": shape.text if shape.has_text_frame else None,
                    "paragraphs": _shape_text_detail(shape, address),
                    "chart": _chart_detail(shape),
                    "estimated_capacity": {
                        "lines_at_18pt": max(1, int(height_inches * 72 / (18 * 1.25))),
                        "characters_per_line_at_18pt": max(
                            1, int(width_inches * 72 / (18 * 0.55))
                        ),
                    },
                }
            )
        slides.append({"address": slide_address, "shapes": shapes})
    return {
        "schema": "ppt-template-detail/v1",
        "source": str(source),
        "slide_size": {
            "width_inches": round(prs.slide_width / EMU_PER_INCH, 3),
            "height_inches": round(prs.slide_height / EMU_PER_INCH, 3),
        },
        "slides": slides,
    }


def _address_parts(
    address: str, *, require_run: bool | None = None
) -> tuple[int, int, int | None, int | None]:
    match = ADDRESS_RE.fullmatch(address)
    if match is None:
        raise SkillError(f"invalid address: {address!r}")
    paragraph = int(match.group("paragraph")) if match.group("paragraph") else None
    run = int(match.group("run")) if match.group("run") else None
    if require_run is True and (paragraph is None or run is None):
        raise SkillError(f"operation requires a run address: {address}")
    if require_run is False and (paragraph is not None or run is not None):
        raise SkillError(f"operation requires a shape address: {address}")
    return int(match.group("slide")), int(match.group("shape")), paragraph, run


def _shape_at(prs: Presentation, slide_index: int, shape_index: int) -> Any:
    if slide_index < 1 or slide_index > len(prs.slides):
        raise SkillError(f"slide address is out of range: {slide_index}")
    shapes = prs.slides[slide_index - 1].shapes
    if shape_index < 1 or shape_index > len(shapes):
        raise SkillError(f"shape address is out of range: slide/{slide_index}/shape/{shape_index}")
    return shapes[shape_index - 1]


def _replace_in_paragraph(paragraph: Any, old: str, new: str) -> int:
    if not paragraph.runs:
        return 0
    replacements = 0
    search_from = 0
    while True:
        texts = [run.text for run in paragraph.runs]
        combined = "".join(texts)
        start = combined.find(old, search_from)
        if start < 0:
            break
        end = start + len(old)
        offsets = []
        cursor = 0
        for text in texts:
            offsets.append((cursor, cursor + len(text)))
            cursor += len(text)
        first = next(index for index, (_, finish) in enumerate(offsets) if finish > start)
        last = next(index for index, (_, finish) in enumerate(offsets) if finish >= end)
        first_start, _ = offsets[first]
        last_start, _ = offsets[last]
        prefix = texts[first][: start - first_start]
        suffix = texts[last][end - last_start :]
        if first == last:
            paragraph.runs[first].text = prefix + new + suffix
        else:
            paragraph.runs[first].text = prefix + new
            for index in range(first + 1, last):
                paragraph.runs[index].text = ""
            paragraph.runs[last].text = suffix
        replacements += 1
        # Skip the replacement itself so values such as "x" -> "xx" cannot loop forever.
        search_from = start + len(new)
    return replacements


def _validate_edit_spec(payload: Any) -> list[dict[str, Any]]:
    if not isinstance(payload, dict):
        raise SkillError("edit spec must be a JSON object")
    schema = payload.get("schema")
    if schema != "ppt-edit-spec/v1":
        raise SkillError(
            f"unsupported edit schema {schema!r}; this CLI accepts only ppt-edit-spec/v1"
        )
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        raise SkillError("edit spec operations must be a non-empty array")
    for operation in operations:
        if not isinstance(operation, dict):
            raise SkillError("each edit operation must be an object")
        if operation.get("op") not in {"replace_text", "set_text", "replace_chart_data"}:
            raise SkillError(f"unknown edit operation: {operation.get('op')!r}")
        if not isinstance(operation.get("address"), str):
            raise SkillError("each edit operation requires an address")
    return operations


def apply_template(source: Path, spec_path: Path, output: Path) -> dict[str, Any]:
    try:
        payload = json.loads(spec_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise SkillError(f"invalid edit spec JSON: {exc}") from exc
    operations = _validate_edit_spec(payload)
    prs = Presentation(source)
    applied = []
    for operation in operations:
        kind = operation["op"]
        address = operation["address"]
        slide_index, shape_index, paragraph_index, run_index = _address_parts(
            address, require_run=kind == "set_text"
        )
        shape = _shape_at(prs, slide_index, shape_index)
        if kind == "replace_text":
            if not shape.has_text_frame:
                raise SkillError(f"replace_text target has no text frame: {address}")
            old = operation.get("old")
            new = operation.get("new")
            if not isinstance(old, str) or not old:
                raise SkillError("replace_text requires a non-empty old string")
            if not isinstance(new, str):
                raise SkillError("replace_text requires a string new value")
            count = sum(
                _replace_in_paragraph(paragraph, old, new)
                for paragraph in shape.text_frame.paragraphs
            )
            if count == 0:
                raise SkillError(f"text {old!r} was not found at {address}")
            applied.append({"op": kind, "address": address, "replacements": count})
        elif kind == "set_text":
            if not shape.has_text_frame or paragraph_index is None or run_index is None:
                raise SkillError(f"set_text target has no text frame: {address}")
            paragraphs = shape.text_frame.paragraphs
            if paragraph_index < 1 or paragraph_index > len(paragraphs):
                raise SkillError(f"paragraph address is out of range: {address}")
            runs = paragraphs[paragraph_index - 1].runs
            if run_index < 1 or run_index > len(runs):
                raise SkillError(f"run address is out of range: {address}")
            text = operation.get("text")
            if not isinstance(text, str):
                raise SkillError("set_text requires a string text value")
            runs[run_index - 1].text = text
            applied.append({"op": kind, "address": address})
        else:
            if not shape.has_chart:
                raise SkillError(f"replace_chart_data target is not a chart: {address}")
            data = _validate_chart(
                {
                    "type": "column",
                    "categories": operation.get("categories"),
                    "series": operation.get("series"),
                }
            )
            shape.chart.replace_data(_chart_data(data))
            applied.append({"op": kind, "address": address})
    _save_presentation(prs, output)
    return {
        "schema": "ppt-edit-result/v1",
        "source": str(source),
        "spec": str(spec_path),
        "output": str(output),
        "applied": applied,
    }


def _issue(code: str, message: str, *, address: str | None = None) -> dict[str, Any]:
    result = {"code": code, "message": message}
    if address:
        result["address"] = address
    return result


def _title_shape(slide: Any) -> Any | None:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "AeloonTitle":
            return shape
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            if shape.placeholder_format.type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
                return shape
    candidates = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() and shape.top < Inches(1.75)
    ]
    return min(candidates, key=lambda item: item.top, default=None)


def _effective_font_size(shape: Any, default: float = 18.0) -> float:
    if not shape.has_text_frame:
        return default
    sizes = [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    return max(sizes, default=default)


def _estimated_overflow(shape: Any) -> tuple[int, int]:
    font_size = _effective_font_size(shape)
    width_points = max(1.0, shape.width / EMU_PER_INCH * 72)
    height_points = max(1.0, shape.height / EMU_PER_INCH * 72)
    chars_per_line = max(1, int(width_points / (font_size * 0.55)))
    required_lines = 0
    for paragraph in shape.text_frame.paragraphs:
        logical_lines = paragraph.text.splitlines() or [""]
        required_lines += sum(
            max(1, math.ceil(len(line) / chars_per_line)) for line in logical_lines
        )
    available_lines = max(1, int(height_points / (font_size * 1.22)))
    return required_lines, available_lines


def _explicit_rgb(color: Any) -> tuple[int, int, int] | None:
    try:
        rgb = color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return tuple(rgb) if rgb is not None else None


def _background_rgb(slide: Any, shape: Any) -> tuple[int, int, int]:
    try:
        shape_color = _explicit_rgb(shape.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        shape_color = None
    if shape_color is not None:
        return shape_color
    try:
        slide_color = _explicit_rgb(slide.background.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        slide_color = None
    return slide_color or (255, 255, 255)


def _relative_luminance(rgb: tuple[int, int, int]) -> float:
    channels = []
    for value in rgb:
        channel = value / 255
        adjusted = (
            channel / 12.92
            if channel <= 0.04045
            else ((channel + 0.055) / 1.055) ** 2.4
        )
        channels.append(adjusted)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_ratio(first: tuple[int, int, int], second: tuple[int, int, int]) -> float:
    first_luminance = _relative_luminance(first)
    second_luminance = _relative_luminance(second)
    lighter = max(first_luminance, second_luminance)
    darker = min(first_luminance, second_luminance)
    return (lighter + 0.05) / (darker + 0.05)


def _content_shape(shape: Any) -> bool:
    return bool(
        (shape.has_text_frame and shape.text.strip())
        or shape.has_chart
        or shape.has_table
        or int(shape.shape_type) == 13
    )


def _overlap_ratio(first: Any, second: Any) -> float:
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    if right <= left or bottom <= top:
        return 0.0
    smaller_area = min(first.width * first.height, second.width * second.height)
    return (right - left) * (bottom - top) / smaller_area if smaller_area else 0.0


def validate_deck(source: Path) -> dict[str, Any]:
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    try:
        with ZipFile(source) as package:
            bad_member = package.testzip()
            if bad_member:
                errors.append(_issue("corrupt_member", f"corrupt package member: {bad_member}"))
            required = {"[Content_Types].xml", "ppt/presentation.xml"}
            missing = sorted(required.difference(package.namelist()))
            if missing:
                errors.append(
                    _issue("missing_package_part", f"missing parts: {', '.join(missing)}")
                )
        prs = Presentation(source)
    except (BadZipFile, KeyError, ValueError, OSError) as exc:
        errors.append(_issue("invalid_pptx", str(exc)))
        return {
            "schema": "ppt-validation/v1",
            "source": str(source),
            "valid": False,
            "errors": errors,
            "warnings": warnings,
            "metrics": {"slides": 0, "shapes": 0, "charts": 0, "tables": 0},
            "visual_qa": "not_run",
        }

    title_sizes = []
    shape_count = chart_count = table_count = 0
    for slide_index, slide in enumerate(prs.slides, start=1):
        title = _title_shape(slide)
        if title is None or not title.text.strip():
            warnings.append(
                _issue(
                    "missing_title",
                    "slide has no detectable title",
                    address=f"slide/{slide_index}",
                )
            )
        elif slide_index > 1:
            title_sizes.append((slide_index, _effective_font_size(title)))
        content_shapes: list[tuple[int, Any]] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_count += 1
            chart_count += int(shape.has_chart)
            table_count += int(shape.has_table)
            address = f"slide/{slide_index}/shape/{shape_index}"
            if _content_shape(shape):
                content_shapes.append((shape_index, shape))
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > prs.slide_width
                or shape.top + shape.height > prs.slide_height
            ):
                warnings.append(
                    _issue(
                        "shape_out_of_bounds",
                        "shape extends outside the slide",
                        address=address,
                    )
                )
            if shape.has_text_frame and shape.text.strip():
                match = PLACEHOLDER_RE.search(shape.text)
                if match:
                    warnings.append(
                        _issue(
                            "placeholder_text",
                            f"possible placeholder text: {match.group(0)!r}",
                            address=address,
                        )
                    )
                required, available = _estimated_overflow(shape)
                if required > available:
                    warnings.append(
                        _issue(
                            "possible_text_overflow",
                            (
                                f"estimated {required} lines for capacity {available}; "
                                "render to verify"
                            ),
                            address=address,
                        )
                    )
                background = _background_rgb(slide, shape)
                low_contrast = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        foreground = _explicit_rgb(run.font.color)
                        if foreground is not None and _contrast_ratio(foreground, background) < 3:
                            low_contrast.append(_contrast_ratio(foreground, background))
                if low_contrast:
                    warnings.append(
                        _issue(
                            "low_text_contrast",
                            f"minimum explicit text contrast is {min(low_contrast):.2f}:1",
                            address=address,
                        )
                    )
        for position, (first_index, first) in enumerate(content_shapes):
            for second_index, second in content_shapes[position + 1 :]:
                ratio = _overlap_ratio(first, second)
                if ratio >= 0.15:
                    warnings.append(
                        _issue(
                            "possible_shape_overlap",
                            (
                                f"content shapes overlap by {ratio:.0%} of the smaller shape; "
                                "render to verify"
                            ),
                            address=(
                                f"slide/{slide_index}/shape/{first_index} and "
                                f"slide/{slide_index}/shape/{second_index}"
                            ),
                        )
                    )
    if len(title_sizes) > 1:
        median = sorted(size for _, size in title_sizes)[len(title_sizes) // 2]
        for slide_index, size in title_sizes:
            if abs(size - median) > 2:
                warnings.append(
                    _issue(
                        "inconsistent_title_size",
                        f"title is {size:g}pt; peer title median is {median:g}pt",
                        address=f"slide/{slide_index}",
                    )
                )
    return {
        "schema": "ppt-validation/v1",
        "source": str(source),
        "valid": not errors,
        "errors": errors,
        "warnings": warnings,
        "metrics": {
            "slides": len(prs.slides),
            "shapes": shape_count,
            "charts": chart_count,
            "tables": table_count,
        },
        "visual_qa": "not_run",
    }


def _libreoffice_command() -> str | None:
    for command in ("soffice", "libreoffice"):
        path = shutil.which(command)
        if path:
            return path
    macos = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    return str(macos) if macos.is_file() else None


def _libreoffice_install_hint() -> str:
    if sys.platform == "darwin":
        return "install LibreOffice with: brew install --cask libreoffice"
    return "install LibreOffice with: sudo apt-get install libreoffice"


def render_deck(
    source: Path, output_dir: Path, *, dpi: int, overwrite: bool
) -> dict[str, Any]:
    command = _libreoffice_command()
    if command is None:
        raise SkillError(
            "LibreOffice is required for PPTX visual rendering; " + _libreoffice_install_hint()
        )
    if dpi < 36 or dpi > 600:
        raise SkillError("dpi must be between 36 and 600")
    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="aeloon-pptx-render-") as temporary:
        temporary_dir = Path(temporary)
        process = subprocess.run(
            [
                command,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(temporary_dir),
                str(source),
            ],
            capture_output=True,
            text=True,
            timeout=180,
            check=False,
        )
        pdf = temporary_dir / f"{source.stem}.pdf"
        if process.returncode != 0 or not pdf.is_file():
            detail = (process.stderr or process.stdout).strip()
            raise SkillError(
                f"LibreOffice failed to render the deck: {detail or 'no PDF produced'}"
            )
        try:
            import pypdfium2 as pdfium
        except ImportError as exc:
            raise SkillError("pypdfium2 is required to render the LibreOffice PDF") from exc
        document = pdfium.PdfDocument(pdf)
        images = []
        try:
            for page_index in range(len(document)):
                target = output_dir / f"slide-{page_index + 1:03d}.png"
                if target.exists() and not overwrite:
                    raise SkillError(f"render output already exists: {target}; pass --overwrite")
                page = document[page_index]
                try:
                    bitmap = page.render(scale=dpi / 72)
                    bitmap.to_pil().save(target, format="PNG")
                    images.append(str(target.resolve()))
                finally:
                    page.close()
        finally:
            document.close()
    return {
        "schema": "ppt-render-result/v1",
        "source": str(source),
        "output_dir": str(output_dir.resolve()),
        "dpi": dpi,
        "slides": len(images),
        "images": images,
        "visual_qa": "rendered_not_inspected",
    }


def _parser(action: str) -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog=f"powerpoint-pptx {action}")
    if action == "build":
        parser.add_argument("source")
        parser.add_argument("output")
        parser.add_argument("--title")
        parser.add_argument("--author")
        parser.add_argument("--accent", default=THEME["accent"])
        parser.add_argument("--font", default="Arial")
        parser.add_argument("--footer", default="")
        parser.add_argument("--overwrite", action="store_true")
    elif action == "inspect-template":
        parser.add_argument("input")
        parser.add_argument("--output")
        parser.add_argument("--overwrite", action="store_true")
    elif action == "apply-template":
        parser.add_argument("input")
        parser.add_argument("spec")
        parser.add_argument("output")
        parser.add_argument("--overwrite", action="store_true")
    elif action == "validate":
        parser.add_argument("input")
        parser.add_argument("--output")
        parser.add_argument("--overwrite", action="store_true")
        parser.add_argument("--strict", action="store_true")
    elif action == "render":
        parser.add_argument("input")
        parser.add_argument("--output-dir", required=True)
        parser.add_argument("--dpi", type=int, default=144)
        parser.add_argument("--overwrite", action="store_true")
    return parser


def main(argv: list[str] | None = None) -> int:
    arguments = list(sys.argv[1:] if argv is None else argv)
    if not arguments or arguments[0] not in VALID_ACTIONS:
        received = arguments[0] if arguments else "<missing>"
        print(
            json.dumps(
                {
                    "error": f"unknown action {received!r}",
                    "valid_actions": list(VALID_ACTIONS),
                },
                ensure_ascii=False,
            ),
            file=sys.stderr,
        )
        return 2
    action = arguments.pop(0)
    args = _parser(action).parse_args(arguments)
    try:
        if action == "build":
            source = _local_file(args.source, suffixes={".md", ".markdown"})
            output = _output_file(
                args.output, source=source, overwrite=args.overwrite, suffix=".pptx"
            )
            result = build_deck(
                source,
                output,
                title=args.title,
                author=args.author,
                accent=args.accent,
                font=args.font,
                footer=args.footer,
            )
        elif action == "inspect-template":
            source = _local_file(args.input, suffixes={".pptx"})
            result = inspect_template(source)
            _write_json(
                Path(args.output) if args.output else None,
                result,
                overwrite=args.overwrite,
            )
        elif action == "apply-template":
            source = _local_file(args.input, suffixes={".pptx"})
            spec = _local_file(args.spec, suffixes={".json"})
            output = _output_file(
                args.output, source=source, overwrite=args.overwrite, suffix=".pptx"
            )
            result = apply_template(source, spec, output)
        elif action == "validate":
            source = _local_file(args.input, suffixes={".pptx"})
            result = validate_deck(source)
            _write_json(
                Path(args.output) if args.output else None,
                result,
                overwrite=args.overwrite,
            )
            print(json.dumps(result, ensure_ascii=False, indent=2))
            if not result["valid"] or (args.strict and result["warnings"]):
                return 1
            return 0
        else:
            source = _local_file(args.input, suffixes={".pptx"})
            result = render_deck(
                source,
                Path(args.output_dir).expanduser().resolve(),
                dpi=args.dpi,
                overwrite=args.overwrite,
            )
    except (SkillError, OSError, ValueError) as exc:
        print(json.dumps({"error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 2
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
