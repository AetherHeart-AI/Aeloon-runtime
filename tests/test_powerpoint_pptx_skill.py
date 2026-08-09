from __future__ import annotations

import hashlib
import importlib.util
import json
import sys
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
import yaml
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor

RESOURCE_ROOT = (
    Path(__file__).parents[1]
    / "aeloon_core"
    / "resources"
    / "skills"
    / "powerpoint-pptx"
)
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def load_cli():
    path = RESOURCE_ROOT / "scripts" / "cli.py"
    spec = importlib.util.spec_from_file_location("test_powerpoint_pptx_cli", path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def build_fixture(tmp_path: Path):
    cli = load_cli()
    image = tmp_path / "metric.png"
    Image.new("RGB", (640, 360), (40, 120, 208)).save(image)
    source = tmp_path / "deck.md"
    source.write_text(
        f"""# Product review

Decisions for the next quarter

```notes
Open with the customer outcome.
```

---

## Adoption accelerated

- Active teams grew 28%
- Retention improved

```chart
{{"type":"column","categories":["Q1","Q2"],"series":[{{"name":"Teams","values":[120,154]}}]}}
```

---

## Segment performance is balanced

| Segment | Growth |
| --- | ---: |
| Enterprise | 31% |
| Growth | 24% |

---

## The new workflow is ready

![Product metric]({image.name})
""",
        encoding="utf-8",
    )
    output = tmp_path / "deck.pptx"
    result = cli.build_deck(
        source,
        output,
        title=None,
        author="Aeloon",
        accent="2878D0",
        font="Arial",
        footer="Internal",
    )
    return cli, source, output, result


def inject_animation_timing(source: Path) -> None:
    with ZipFile(source) as package:
        entries = {name: package.read(name) for name in package.namelist()}
    slide_name = "ppt/slides/slide2.xml"
    root = etree.fromstring(entries[slide_name])
    timing = etree.fromstring(
        (
            f'<p:timing xmlns:p="{PML_NS}"><p:tnLst><p:par>'
            '<p:cTn id="42" dur="indefinite" restart="never"/>'
            "</p:par></p:tnLst></p:timing>"
        ).encode()
    )
    root.append(timing)
    entries[slide_name] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")
    rebuilt = source.with_suffix(".rebuilt.pptx")
    with ZipFile(rebuilt, "w", compression=ZIP_DEFLATED) as package:
        for name, data in entries.items():
            package.writestr(name, data)
    rebuilt.replace(source)


def preserved_parts(source: Path) -> dict[str, bytes]:
    prefixes = ("ppt/media/", "ppt/theme/", "ppt/slideMasters/", "ppt/slideLayouts/")
    with ZipFile(source) as package:
        return {
            name: package.read(name)
            for name in package.namelist()
            if name.startswith(prefixes) or name == "ppt/slides/slide4.xml"
        }


def test_skill_metadata_specs_and_license_are_complete() -> None:
    skill_text = (RESOURCE_ROOT / "SKILL.md").read_text(encoding="utf-8")
    _, frontmatter, body = skill_text.split("---", 2)
    metadata = yaml.safe_load(frontmatter)
    assert metadata["name"] == "powerpoint-pptx"
    assert set(metadata) == {"name", "description"}
    assert "present_files" in body
    for action in ("build", "inspect-template", "apply-template", "validate", "render"):
        assert action in body
    interface = yaml.safe_load(
        (RESOURCE_ROOT / "agents" / "openai.yaml").read_text(encoding="utf-8")
    )["interface"]
    assert "$powerpoint-pptx" in interface["default_prompt"]
    license_text = (RESOURCE_ROOT / "LICENSE.txt").read_text(encoding="utf-8")
    assert "MIT License" in license_text
    assert "source code only" in license_text
    assert "templates" in license_text
    assert "Aeloon modification and provenance notice" in license_text
    assert not (RESOURCE_ROOT / "templates").exists()


def test_build_creates_editable_slides_notes_table_image_and_chart(tmp_path: Path) -> None:
    cli, _, output, result = build_fixture(tmp_path)

    assert result["slides"] == 4
    assert result["theme"] == "aeloon-minimal-16x9/v1"
    prs = Presentation(output)
    assert len(prs.slides) == 4
    assert prs.slide_width / cli.EMU_PER_INCH == pytest.approx(13.333333, abs=0.001)
    assert prs.slide_height / cli.EMU_PER_INCH == pytest.approx(7.5)
    assert "customer outcome" in prs.slides[0].notes_slide.notes_text_frame.text
    assert any(shape.has_chart for shape in prs.slides[1].shapes)
    assert any(shape.has_table for shape in prs.slides[2].shapes)
    assert any(shape.shape_type == 13 for shape in prs.slides[3].shapes)
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Active teams grew 28%" in all_text

    validation = cli.validate_deck(output)
    assert validation["valid"] is True
    assert validation["metrics"]["charts"] == 1
    assert validation["metrics"]["tables"] == 1
    assert validation["warnings"] == []


def test_inspect_and_apply_template_preserve_input_and_run_format(tmp_path: Path) -> None:
    cli, _, source, _ = build_fixture(tmp_path)
    inject_animation_timing(source)
    before = hashlib.sha256(source.read_bytes()).hexdigest()
    untouched_before = preserved_parts(source)
    detail = cli.inspect_template(source)
    assert detail["schema"] == "ppt-template-detail/v1"
    assert detail["slide_size"]["height_inches"] == 7.5

    title_shape = next(
        shape
        for shape in detail["slides"][1]["shapes"]
        if shape["name"] == "AeloonTitle"
    )
    chart_shape = next(
        shape for shape in detail["slides"][1]["shapes"] if shape["chart"] is not None
    )
    title_run = title_shape["paragraphs"][0]["runs"][0]
    original_size = title_run["font"]["size_pt"]
    spec = tmp_path / "edits.json"
    spec.write_text(
        json.dumps(
            {
                "schema": "ppt-edit-spec/v1",
                "operations": [
                    {
                        "op": "replace_text",
                        "address": title_shape["address"],
                        "old": "Adoption accelerated",
                        "new": "Adoption reached a new high",
                    },
                    {
                        "op": "replace_chart_data",
                        "address": chart_shape["address"],
                        "categories": ["Q1", "Q2"],
                        "series": [{"name": "Teams", "values": [130, 180]}],
                    },
                ],
            }
        ),
        encoding="utf-8",
    )
    output = tmp_path / "edited.pptx"
    result = cli.apply_template(source, spec, output)

    assert hashlib.sha256(source.read_bytes()).hexdigest() == before
    assert preserved_parts(output) == untouched_before
    assert len(result["applied"]) == 2
    with ZipFile(output) as package:
        edited_slide = etree.fromstring(package.read("ppt/slides/slide2.xml"))
    timing = edited_slide.find(f"{{{PML_NS}}}timing")
    assert timing is not None
    assert timing.xpath("string(.//p:cTn/@id)", namespaces={"p": PML_NS}) == "42"
    edited = cli.inspect_template(output)
    edited_title = next(
        shape
        for shape in edited["slides"][1]["shapes"]
        if shape["name"] == "AeloonTitle"
    )
    assert edited_title["text"] == "Adoption reached a new high"
    assert edited_title["paragraphs"][0]["runs"][0]["font"]["size_pt"] == original_size
    edited_chart = next(
        shape for shape in edited["slides"][1]["shapes"] if shape["chart"] is not None
    )
    assert edited_chart["chart"]["series"][0]["values"] == [130.0, 180.0]


def test_validation_reports_placeholder_and_possible_overflow(tmp_path: Path) -> None:
    cli = load_cli()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(0, 0, 3000000, 500000)
    title.name = "AeloonTitle"
    title.text = "Validation"
    body = slide.shapes.add_textbox(100000, 700000, 1000000, 180000)
    body.text = "PLACEHOLDER " + ("dense text " * 100)
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(250, 250, 250)
    body.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(245, 245, 245)
    overlap = slide.shapes.add_textbox(150000, 720000, 1000000, 180000)
    overlap.text = "Overlapping content"
    source = tmp_path / "invalid-content.pptx"
    prs.save(source)

    result = cli.validate_deck(source)
    codes = {warning["code"] for warning in result["warnings"]}
    assert result["valid"] is True
    assert "placeholder_text" in codes
    assert "possible_text_overflow" in codes
    assert "low_text_contrast" in codes
    assert "possible_shape_overlap" in codes


def test_render_without_libreoffice_has_actionable_install_hint(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    cli, _, source, _ = build_fixture(tmp_path)
    monkeypatch.setattr(cli, "_libreoffice_command", lambda: None)
    with pytest.raises(cli.SkillError, match="install LibreOffice"):
        cli.render_deck(source, tmp_path / "rendered", dpi=144, overwrite=False)


def test_unknown_action_lists_valid_actions(capsys: pytest.CaptureFixture[str]) -> None:
    cli = load_cli()
    assert cli.main(["merge"]) == 2
    error = json.loads(capsys.readouterr().err)
    assert error["valid_actions"] == list(cli.VALID_ACTIONS)
